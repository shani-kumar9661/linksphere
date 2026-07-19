import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def main():
    prs = Presentation()
    # Set to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color palette
    DARK_BG = RGBColor(15, 23, 42)      # Slate 900
    LIGHT_TEXT = RGBColor(241, 245, 249) # Slate 100
    ACCENT_BLUE = RGBColor(56, 189, 248) # Cyan 400
    MUTED_TEXT = RGBColor(148, 163, 184) # Slate 400
    RED_ACCENT = RGBColor(239, 68, 68)   # Red 500

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

    def add_header(slide, title_text, category_text="LINKSphere PLATFORM"):
        # Add category tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = 'Arial'
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_BLUE
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.font.color.rgb = LIGHT_TEXT

    blank_layout = prs.slide_layouts[6]

    # ----------------------------------------------------
    # SLIDE 1: Title Slide
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True

    p_logo = tf.paragraphs[0]
    p_logo.text = "🔗 LINKSphere"
    p_logo.font.name = 'Arial'
    p_logo.font.size = Pt(64)
    p_logo.font.bold = True
    p_logo.font.color.rgb = ACCENT_BLUE
    p_logo.alignment = PP_ALIGN.CENTER
    p_logo.space_after = Pt(10)

    p_sub = tf.add_paragraph()
    p_sub.text = "Cloud-Native Smart URL Management Platform"
    p_sub.font.name = 'Arial'
    p_sub.font.size = Pt(24)
    p_sub.font.bold = True
    p_sub.font.color.rgb = LIGHT_TEXT
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.space_after = Pt(20)

    p_pres = tf.add_paragraph()
    p_pres.text = "Advanced Full-Stack Engineering College Project Presentation"
    p_pres.font.name = 'Arial'
    p_pres.font.size = Pt(16)
    p_pres.font.italic = True
    p_pres.font.color.rgb = MUTED_TEXT
    p_pres.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # SLIDE 2: Problem vs Solution
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "The Problem & The Solution")

    # Left Column (Challenge)
    left_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    p_lh = tf_left.paragraphs[0]
    p_lh.text = "❌ THE CHALLENGE"
    p_lh.font.name = 'Arial'
    p_lh.font.size = Pt(20)
    p_lh.font.bold = True
    p_lh.font.color.rgb = RED_ACCENT
    p_lh.space_after = Pt(14)

    bullets_left = [
        "Long, messy URLs are difficult to share, print, or scan.",
        "Generic shorteners lack analytical insights and visitor details.",
        "No way to set expiration dates or password lock sensitive links.",
        "Scaling databases to handle millions of redirects causes high latency."
    ]
    for b in bullets_left:
        p = tf_left.add_paragraph()
        p.text = "•  " + b
        p.font.name = 'Arial'
        p.font.size = Pt(16)
        p.font.color.rgb = MUTED_TEXT
        p.space_after = Pt(12)

    # Right Column (Solution)
    right_box = slide2.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    p_rh = tf_right.paragraphs[0]
    p_rh.text = "🛡️ THE SOLUTION: LINKSphere"
    p_rh.font.name = 'Arial'
    p_rh.font.size = Pt(20)
    p_rh.font.bold = True
    p_rh.font.color.rgb = ACCENT_BLUE
    p_rh.space_after = Pt(14)

    bullets_right = [
        "Sleek and clean URL shortening with custom aliases.",
        "Dynamic QR code generation ready to download as PNG.",
        "Link Access Controls: Secure password protection & automated expiry.",
        "High-Performance Caching: Redis-backed redirection in under 2ms.",
        "Real-time visual dashboard for device, browser, geo, and traffic stats."
    ]
    for b in bullets_right:
        p = tf_right.add_paragraph()
        p.text = "•  " + b
        p.font.name = 'Arial'
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_TEXT
        p.space_after = Pt(12)


    # ----------------------------------------------------
    # SLIDE 3: System Architecture
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "System Architecture & Tier Flow")

    tiers = [
        {"title": "1. FRONTEND TIER", "tech": "React + Vite + Tailwind", "desc": "Handles interactive dashboard rendering, Recharts visualization, dynamic QR downloads, and authentication forms."},
        {"title": "2. GATEWAY PROXY", "tech": "Vercel Engine / Nginx", "desc": "Manages route rewriting, proxying client API calls to backend without CORS errors, and securing SSL/TLS termination."},
        {"title": "3. BACKEND API TIER", "tech": "Node.js + Express.js", "desc": "MVC architecture processing link creation, device/browser metadata parsing, JWT verification, and background jobs."},
        {"title": "4. DATABASE & CACHE", "tech": "MongoDB + Upstash Redis", "desc": "MongoDB stores persistent data (Users, Links, Clicks). Redis acts as memory cache for hot redirection and rate-limiter storage."}
    ]

    for i, t in enumerate(tiers):
        x = Inches(0.8 + i * 2.95)
        box = slide3.shapes.add_textbox(x, Inches(1.8), Inches(2.8), Inches(4.8))
        tf = box.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = t["title"]
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = ACCENT_BLUE
        p_t.space_after = Pt(8)
        
        p_tech = tf.add_paragraph()
        p_tech.text = t["tech"]
        p_tech.font.name = 'Arial'
        p_tech.font.size = Pt(14)
        p_tech.font.italic = True
        p_tech.font.bold = True
        p_tech.font.color.rgb = LIGHT_TEXT
        p_tech.space_after = Pt(12)
        
        p_desc = tf.add_paragraph()
        p_desc.text = t["desc"]
        p_desc.font.name = 'Arial'
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = MUTED_TEXT
        p_desc.space_after = Pt(8)


    # ----------------------------------------------------
    # SLIDE 4: Key Smart Features
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "Smart Features & Link Controls")

    features = [
        {"icon": "🔒", "title": "Password Protection", "desc": "URLs can be locked with a password. Visitors must verify the correct password before being redirected to the target page."},
        {"icon": "📅", "title": "Link Expiration", "desc": "Set auto-expiry (1 day, 7 days, 30 days, or custom date/time) after which links automatically deactivate."},
        {"icon": "📱", "title": "Dynamic QR Codes", "desc": "Generates scannable QR codes instantly for every shortlink, available for download as high-res PNG format."},
        {"icon": "🏷️", "title": "Categorization & Tags", "desc": "Organize links by adding custom tags (e.g. marketing, social) and categories for simple search & filtering."}
    ]

    for i, f in enumerate(features):
        row = i // 2
        col = i % 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(1.8 + row * 2.5)
        
        box = slide4.shapes.add_textbox(x, y, Inches(5.6), Inches(2.2))
        tf = box.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = f"{f['icon']}  {f['title']}"
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = ACCENT_BLUE
        p_t.space_after = Pt(8)
        
        p_d = tf.add_paragraph()
        p_d.text = f["desc"]
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(15)
        p_d.font.color.rgb = LIGHT_TEXT
        p_d.space_after = Pt(8)


    # ----------------------------------------------------
    # SLIDE 5: Real-Time Analytics
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "Real-Time Analytics & Visitor Tracking")

    # Left Side: Graphic description
    desc_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    
    p_dh = tf_desc.paragraphs[0]
    p_dh.text = "📈 VISITOR METADATA ANALYSIS"
    p_dh.font.name = 'Arial'
    p_dh.font.size = Pt(20)
    p_dh.font.bold = True
    p_dh.font.color.rgb = ACCENT_BLUE
    p_dh.space_after = Pt(14)

    analytics_bullets = [
        "User-Agent Parsing: Auto-extracts Operating System (Windows, macOS, Linux), Device type (Desktop, Mobile, Tablet), and Browser type (Chrome, Firefox, Safari).",
        "Referrer Tracking: Identifies traffic sources to reveal whether visitors came directly or via social media platforms (LinkedIn, Twitter).",
        "GDPR-Compliant Geolocation: Logs unique visits and locations while hashes IP addresses to ensure data privacy."
    ]
    for b in analytics_bullets:
        p = tf_desc.add_paragraph()
        p.text = "•  " + b
        p.font.name = 'Arial'
        p.font.size = Pt(15)
        p.font.color.rgb = LIGHT_TEXT
        p.space_after = Pt(12)

    # Right Side: Key KPIs
    kpi_box = slide5.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_kpi = kpi_box.text_frame
    tf_kpi.word_wrap = True
    
    p_kh = tf_kpi.paragraphs[0]
    p_kh.text = "🎯 KEY PERFORMANCE METRICS"
    p_kh.font.name = 'Arial'
    p_kh.font.size = Pt(20)
    p_kh.font.bold = True
    p_kh.font.color.rgb = LIGHT_TEXT
    p_kh.space_after = Pt(14)

    kpi_bullets = [
        "Click Timelines: Interactive line charts representing daily, weekly, or monthly click frequencies.",
        "Browser & OS Distribution: Visual pie charts enabling administrative review of user demography.",
        "Link Performance Score: Shows active vs. archived links to evaluate marketing campaigns."
    ]
    for b in kpi_bullets:
        p = tf_kpi.add_paragraph()
        p.text = "•  " + b
        p.font.name = 'Arial'
        p.font.size = Pt(15)
        p.font.color.rgb = MUTED_TEXT
        p.space_after = Pt(12)


    # ----------------------------------------------------
    # SLIDE 6: Performance & Security
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "Performance Caching & Security")

    # Left box: Redis caching
    cache_box = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_cache = cache_box.text_frame
    tf_cache.word_wrap = True
    
    p_ch = tf_cache.paragraphs[0]
    p_ch.text = "⚡ SUB-2MS CACHING VIA REDIS"
    p_ch.font.name = 'Arial'
    p_ch.font.size = Pt(20)
    p_ch.font.bold = True
    p_ch.font.color.rgb = ACCENT_BLUE
    p_ch.space_after = Pt(14)

    cache_bullets = [
        "Database Bypass: Popular redirected links are cached directly in Redis memory, bypassing slow disk-based MongoDB queries.",
        "Rate-Limiting Storage: Redis stores rate-limiting counters to verify user request frequencies without overloading the DB.",
        "Automatic Invalidation: Cache entries auto-expire and update when URLs are modified or deleted by creators."
    ]
    for b in cache_bullets:
        p = tf_cache.add_paragraph()
        p.text = "•  " + b
        p.font.name = 'Arial'
        p.font.size = Pt(15)
        p.font.color.rgb = LIGHT_TEXT
        p.space_after = Pt(12)

    # Right box: Security
    sec_box = slide6.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_sec = sec_box.text_frame
    tf_sec.word_wrap = True
    
    p_sh = tf_sec.paragraphs[0]
    p_sh.text = "🔒 ENTERPRISE-GRADE SECURITY"
    p_sh.font.name = 'Arial'
    p_sh.font.size = Pt(20)
    p_sh.font.bold = True
    p_sh.font.color.rgb = LIGHT_TEXT
    p_sh.space_after = Pt(14)

    sec_bullets = [
        "JWT Token Rotation: Short-lived access tokens (15 mins) combined with httpOnly secure cookies for refresh tokens prevent token hijacking.",
        "Injection Safeguards: Integrated sanitization middleware blocks NoSQL injection attacks and Cross-Site Scripting (XSS).",
        "Brute-Force Lockout: Limits login API requests to a maximum of 10 requests per 15 minutes per IP address."
    ]
    for b in sec_bullets:
        p = tf_sec.add_paragraph()
        p.text = "•  " + b
        p.font.name = 'Arial'
        p.font.size = Pt(15)
        p.font.color.rgb = MUTED_TEXT
        p.space_after = Pt(12)


    # ----------------------------------------------------
    # SLIDE 7: Summary & Conclusion
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "Deployment & Key Deliverables")

    delivs = [
        {"title": "100% Cloud-Native", "desc": "Hosted on serverless infrastructure. Vercel for UI, Render for Server, MongoDB Atlas and Upstash for services."},
        {"title": "No-Credit-Card Stack", "desc": "Implemented using exclusively free cloud tiering, demonstrating cost-effective enterprise architecture scaling."},
        {"title": "Fully Decoupled Build", "desc": "Frontend and Backend are independent, communicating securely through Vercel's edge-level proxy rewrites."}
    ]

    for i, d in enumerate(delivs):
        x = Inches(0.8 + i * 3.9)
        box = slide7.shapes.add_textbox(x, Inches(2.2), Inches(3.7), Inches(4.0))
        tf = box.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = d["title"]
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = ACCENT_BLUE
        p_t.space_after = Pt(12)
        
        p_d = tf.add_paragraph()
        p_d.text = d["desc"]
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(16)
        p_d.font.color.rgb = LIGHT_TEXT
        p_d.space_after = Pt(8)

    prs.save("LinkSphere_Presentation.pptx")
    print("SUCCESS: LinkSphere_Presentation.pptx generated successfully!")

if __name__ == "__main__":
    main()
